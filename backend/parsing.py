import io
import re
from pypdf import PdfReader
from pptx import Presentation

def extract_text(filename: str, raw_bytes: bytes) -> str:
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        reader = PdfReader(io.BytesIO(raw_bytes))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if ext in ("ppt", "pptx"):
        prs = Presentation(io.BytesIO(raw_bytes))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_parts.append(shape.text_frame.text)
        return "\n".join(text_parts)
    return raw_bytes.decode("utf-8", errors="ignore")

def chunk_by_heading(text: str) -> dict:
    """Naive split into chapters/sections. Falls back to one 'Full Document' chunk."""
    lines = text.split("\n")
    chapters = {}
    current_title = "Full Document"
    current_lines = []
    heading_pattern = re.compile(r"^(chapter|section|unit)\s+\d+", re.IGNORECASE)
    for line in lines:
        stripped = line.strip()
        if heading_pattern.match(stripped) and len(stripped) < 80:
            if current_lines:
                chapters[current_title] = "\n".join(current_lines)
            current_title = stripped
            current_lines = []
        else:
            current_lines.append(line)
    chapters[current_title] = "\n".join(current_lines)
    return chapters